#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Doc-to-Content v10

Что добавлено поверх v9:
- CLI режим без UI
- пакетная обработка через командную строку
- zip-архив результата
- единое ядро для desktop / bot / automation
- удобно под Telegram-бота и серверную обёртку

Примеры запуска:
    py pdf_md_gui.py                  # GUI
    py pdf_md_gui.py --cli --input fail1.pdf --client demo --preset single --result summary --mode AI
    py pdf_md_gui.py --cli --input ./docs --client demo --preset creator --mode AI --zip

Установка:
    py -m pip install customtkinter pymupdf python-docx python-pptx beautifulsoup4 markdownify requests
"""

from __future__ import annotations

import argparse
import json
import queue
import re
import shutil
import threading
import zipfile
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

import customtkinter as ctk
from tkinter import filedialog, messagebox

import fitz
import requests
from bs4 import BeautifulSoup
from docx import Document
from markdownify import markdownify as html_to_markdown
from pptx import Presentation


# -------------------- THEME --------------------

ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")

APP_BG = "#0b1020"
CARD_BG = "#121a2b"
CARD_ALT = "#0f1727"
BORDER = "#22314f"
TEXT = "#eef4ff"
TEXT_DIM = "#93a4c3"
ACCENT = "#4f8cff"
ACCENT_HOVER = "#3d74d4"

SUPPORTED_INPUTS = {".pdf", ".docx", ".pptx", ".html", ".htm"}
SETTINGS_PATH = Path("doc_to_content_settings.json")
HISTORY_PATH = Path("doc_to_content_history.jsonl")

FREE_MODELS = [
    "openrouter/free",
    "mistralai/mistral-7b-instruct:free",
    "openchat/openchat-7b:free",
    "nousresearch/nous-capybara-7b:free",
]

PROMPTS = {
    "text": """
Очисти текст от мусора верстки, склей сломанные слова и переносы, сохрани смысл без фантазий.
Исправь очевидные языковые ошибки, если они появились из-за парсинга.
Верни только чистый читабельный текст.
""".strip(),
    "markdown": """
Преврати сырой текст в чистый структурированный Markdown.
Требования:
- восстанови заголовки, абзацы и списки
- убери артефакты парсинга и пустые повторы
- не выдумывай факты
- не добавляй вступлений и комментариев
Верни только итоговый Markdown.
""".strip(),
    "summary": """
Сделай краткое, но профессиональное summary текста.

Требования:
- Пиши как для клиента или бизнес-презентации
- Без слов «Совет:», «Рекомендуется», «Старайтесь»
- Убери общие фразы, банальности и пустые выводы
- Исправь грамматику и формулировки
- Используй единый деловой и уверенный тон

Формат:
- 1–2 абзаца
- Допускается маркированный список, но без «Совет:»

Важно:
- Ничего не выдумывай
- Работай только с исходным текстом
- Верни только итоговый текст
""".strip(),
    "article": """
Преврати текст в полноценную читабельную статью.

Требования:
- Сделай сильный заголовок
- Добавь понятную структуру: заголовки, подзаголовки, абзацы
- Убери мусор парсинга, повторы и канцелярит
- Не добавляй факты, которых нет в исходнике
- Тон: живой, уверенный, профессиональный

Формат:
- Готовая статья в Markdown
- Без пояснений от себя
""".strip(),
    "posts": """
Сделай из текста 5–10 коротких постов для Telegram.

Требования:
- Каждый пост должен быть отдельным и законченным
- Нормальный человеческий русский, без воды
- Без слов «Совет:» и банальных наставлений
- Можно делать разные углы подачи: инсайт, кейс, мысль, позиция
- Не выдумывай факты

Формат:
- Каждый пост отделяй строкой: ---
- Верни только готовые посты
""".strip(),
    "faq": """
Сделай FAQ в формате вопрос-ответ.

Требования:
- Сгруппируй материал логично
- Формулируй вопросы коротко и понятно
- Ответы — по делу, без воды
- Не выдумывай факты

Формат:
- Markdown
- Вопрос как подзаголовок, ниже ответ
- Верни только итоговый FAQ
""".strip(),
}

PACK_PRESETS = {
    "single": [("markdown", ".md")],
    "creator": [("summary", ".md"), ("article", ".md"), ("posts", ".md")],
    "knowledge": [("summary", ".md"), ("faq", ".md"), ("markdown", ".md")],
}


# -------------------- MODELS --------------------

@dataclass
class JobResult:
    source: Path
    target: Path | None
    ok: bool
    message: str


# -------------------- HELPERS --------------------

def now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def slugify(value: str) -> str:
    value = value.strip().lower()
    value = re.sub(r"[^a-zA-Zа-яА-Я0-9_-]+", "_", value)
    value = re.sub(r"_+", "_", value).strip("_")
    return value or "client"


def read_json(path: Path, default):
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return default


def write_json(path: Path, data) -> None:
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def append_jsonl(path: Path, data) -> None:
    with path.open("a", encoding="utf-8") as f:
        f.write(json.dumps(data, ensure_ascii=False) + "\n")


def normalize_text(text: str) -> str:
    text = text.replace("\u00a0", " ")
    text = text.replace("\xad", "")
    text = text.replace("\u2011", "-")
    text = text.replace("\u2013", "-")
    text = text.replace("\u2014", "—")
    text = text.replace("Telegramканала", "Telegram-канала")
    text = text.replace("Telegramканал", "Telegram-канал")
    text = text.replace("Eventкейтеринг", "Event-кейтеринг")
    text = text.replace("￾", "-")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def light_cleanup(text: str) -> str:
    text = normalize_text(text)
    text = re.sub(r"([А-Яа-яA-Za-z])\n([а-яa-z])", r"\1\2", text)
    text = re.sub(r"([А-Яа-яA-Za-z])-\n([А-Яа-яA-Za-z])", r"\1\2", text)

    lines = [line.strip() for line in text.splitlines()]
    out: list[str] = []
    buf = ""

    def flush() -> None:
        nonlocal buf
        if buf:
            out.append(buf.strip())
            buf = ""

    for line in lines:
        if not line:
            flush()
            if out and out[-1] != "":
                out.append("")
            continue

        if re.match(r"^(#|[-*]|\d+[.)])\s+", line):
            flush()
            out.append(line)
            continue

        if not buf:
            buf = line
            continue

        if re.search(r"[.!?:;…]$", buf):
            flush()
            buf = line
        else:
            buf += " " + line

    flush()

    cleaned = "\n".join(out)
    cleaned = re.sub(r"([А-Яа-яA-Za-z])\s+([,:;.!?])", r"\1\2", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip() + "\n"


def _build_headers(api_key: str) -> dict[str, str]:
    headers = {
        "Content-Type": "application/json",
        "HTTP-Referer": "http://localhost",
        "X-Title": "doc-to-content",
    }
    if api_key.strip():
        headers["Authorization"] = f"Bearer {api_key.strip()}"
    return headers


def ai_cleanup(raw_text: str, *, api_url: str, api_key: str, model: str, output_format: str, extra_prompt: str, timeout: int = 180) -> str:
    format_hint = "Markdown" if output_format == ".md" else "plain text"
    prompt = f"""
Ты редактор-конвертер документов.
Тебе дают сырой текст, выдранный из PDF, презентации или документа.

Твоя задача:
1. Убрать мусор верстки.
2. Склеить сломанные слова и переносы.
3. Восстановить читабельную структуру.
4. Не придумывать факты и не дописывать отсебятину.
5. Вернуть результат в формате: {format_hint}.

Специальная задача:
{extra_prompt}

Сырой текст:
{raw_text}
""".strip()

    payload = {
        "model": model.strip(),
        "messages": [
            {"role": "system", "content": "Ты приводишь сырой текст документов в аккуратный и полезный вид без фантазий."},
            {"role": "user", "content": prompt},
        ],
        "temperature": 0.1,
    }

    response = requests.post(api_url.strip(), headers=_build_headers(api_key), json=payload, timeout=timeout)
    response.raise_for_status()
    data = response.json()

    try:
        return data["choices"][0]["message"]["content"].strip() + "\n"
    except Exception as exc:
        raise RuntimeError(f"Не удалось разобрать ответ API: {exc}\nОтвет: {json.dumps(data, ensure_ascii=False)[:800]}")


def ai_cleanup_with_fallback(raw_text: str, *, api_url: str, api_key: str, model: str, output_format: str, extra_prompt: str, timeout: int = 180, log_fn=None) -> str:
    models_to_try = [model.strip()] if model.strip() else []
    for fallback in FREE_MODELS:
        if fallback not in models_to_try:
            models_to_try.append(fallback)

    last_exc = None
    for candidate in models_to_try:
        try:
            if log_fn:
                log_fn(f"Пробую модель: {candidate}")
            return ai_cleanup(
                raw_text,
                api_url=api_url,
                api_key=api_key,
                model=candidate,
                output_format=output_format,
                extra_prompt=extra_prompt,
                timeout=timeout,
            )
        except Exception as exc:
            last_exc = exc
            if log_fn:
                log_fn(f"Модель не взлетела: {candidate} -> {exc}")
            continue

    raise RuntimeError(f"Не удалось получить ответ ни от одной модели. Последняя ошибка: {last_exc}")


def zip_directory(src_dir: Path, zip_path: Path) -> Path:
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for path in src_dir.rglob("*"):
            if path.is_file():
                zf.write(path, path.relative_to(src_dir))
    return zip_path


def collect_input_paths(input_path: Path) -> list[Path]:
    if input_path.is_file():
        return [input_path] if input_path.suffix.lower() in SUPPORTED_INPUTS else []
    if input_path.is_dir():
        return sorted([p for p in input_path.rglob("*") if p.suffix.lower() in SUPPORTED_INPUTS], key=lambda p: p.name.lower())
    return []


# -------------------- EXTRACTORS --------------------

def pdf_to_raw_text(path: Path, as_markdown: bool) -> str:
    doc = fitz.open(path)
    pages: list[str] = []
    for page in doc:
        text = normalize_text(page.get_text("text"))
        if text:
            pages.append(text)
    sep = "\n\n---\n\n" if as_markdown else "\n\n"
    return sep.join(pages).strip() + "\n"


def docx_to_raw_text(path: Path, as_markdown: bool) -> str:
    doc = Document(path)
    parts: list[str] = []
    for p in doc.paragraphs:
        text = normalize_text(p.text)
        if not text:
            parts.append("")
            continue
        style_name = (p.style.name or "").lower() if p.style else ""
        if as_markdown and "heading" in style_name:
            level_match = re.search(r"(\d+)", style_name)
            level = min(int(level_match.group(1)), 6) if level_match else 2
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)
    return "\n".join(parts).strip() + "\n"


def pptx_to_raw_text(path: Path, as_markdown: bool) -> str:
    prs = Presentation(path)
    slides: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        bits: list[str] = []
        if as_markdown:
            bits.append(f"# Слайд {idx}")
            bits.append("")
        for shape in slide.shapes:
            text = normalize_text(getattr(shape, "text", ""))
            if text:
                bits.append(text)
                if as_markdown:
                    bits.append("")
        slide_text = "\n".join(bits).strip()
        if slide_text:
            slides.append(slide_text)
    sep = "\n\n---\n\n" if as_markdown else "\n\n"
    return sep.join(slides).strip() + "\n"


def html_to_raw_text(path: Path, as_markdown: bool) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    if as_markdown:
        return normalize_text(html_to_markdown(raw)) + "\n"
    soup = BeautifulSoup(raw, "html.parser")
    return normalize_text(soup.get_text("\n")) + "\n"


def extract_raw_text(path: Path, output_format: str) -> str:
    ext = path.suffix.lower()
    as_markdown = output_format == ".md"
    if ext == ".pdf":
        return pdf_to_raw_text(path, as_markdown)
    if ext == ".docx":
        return docx_to_raw_text(path, as_markdown)
    if ext == ".pptx":
        return pptx_to_raw_text(path, as_markdown)
    if ext in {".html", ".htm"}:
        return html_to_raw_text(path, as_markdown)
    raise ValueError(f"Формат не поддерживается: {ext}")


# -------------------- CORE SERVICE --------------------

class CoreService:
    def __init__(self, *, process_mode: str, client_name: str, output_folder_name: str, use_output_folder: bool, copy_source: bool, api_url: str, api_key: str, model_name: str, ai_preclean: bool, logger=None):
        self.process_mode = process_mode
        self.client_name = client_name
        self.output_folder_name = output_folder_name
        self.use_output_folder = use_output_folder
        self.copy_source = copy_source
        self.api_url = api_url
        self.api_key = api_key
        self.model_name = model_name
        self.ai_preclean = ai_preclean
        self.logger = logger or (lambda msg: None)

    def _make_client_dir(self, src: Path) -> Path:
        if self.use_output_folder:
            root = src.parent / self.output_folder_name.strip()
            root.mkdir(parents=True, exist_ok=True)
        else:
            root = src.parent
        client_dir = root / f"{slugify(self.client_name)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        client_dir.mkdir(parents=True, exist_ok=True)
        return client_dir

    def _process_single_result(self, raw: str, result_type: str, ext: str) -> str:
        if self.process_mode == "RAW":
            return raw
        if self.process_mode == "CLEAN":
            return light_cleanup(raw)
        if self.process_mode == "AI":
            prepared = light_cleanup(raw) if self.ai_preclean else raw
            prompt = PROMPTS[result_type]
            return ai_cleanup_with_fallback(
                prepared,
                api_url=self.api_url,
                api_key=self.api_key,
                model=self.model_name,
                output_format=ext,
                extra_prompt=prompt,
                log_fn=self.logger,
            )
        raise ValueError(f"Неизвестный режим: {self.process_mode}")

    def process_file(self, src: Path, jobs: list[tuple[str, str]]) -> Path:
        self.logger(f"Старт: {src.name}")
        client_dir = self._make_client_dir(src)
        outputs_dir = client_dir / "outputs"
        outputs_dir.mkdir(parents=True, exist_ok=True)

        if self.copy_source:
            shutil.copy2(src, client_dir / src.name)

        manifest = {
            "timestamp": now_iso(),
            "client_name": self.client_name,
            "source": str(src),
            "process_mode": self.process_mode,
            "jobs": [],
        }

        for result_type, out_ext in jobs:
            raw = extract_raw_text(src, out_ext)
            final = self._process_single_result(raw, result_type, out_ext)
            target = outputs_dir / f"{src.stem}_{result_type}{out_ext}"
            target.write_text(final, encoding="utf-8")
            manifest["jobs"].append({
                "result_type": result_type,
                "path": str(target),
                "ext": out_ext,
            })
            self.logger(f"  -> {target.name}")

        write_json(client_dir / "manifest.json", manifest)
        append_jsonl(HISTORY_PATH, manifest)
        self.logger(f"Готово: {client_dir}")
        return client_dir


# -------------------- GUI --------------------

class App(ctk.CTk):
    def __init__(self) -> None:
        super().__init__()
        self.title("Doc-to-Content v10")
        self.geometry("1360x920")
        self.minsize(1200, 800)
        self.configure(fg_color=APP_BG)

        settings = read_json(SETTINGS_PATH, {})

        self.files: list[Path] = []
        self.queue: queue.Queue[tuple[str, object]] = queue.Queue()

        self.process_mode = ctk.StringVar(value=settings.get("process_mode", "AI"))
        self.single_result_type = ctk.StringVar(value=settings.get("single_result_type", "markdown"))
        self.pack_preset = ctk.StringVar(value=settings.get("pack_preset", "single"))
        self.use_output_folder = ctk.BooleanVar(value=settings.get("use_output_folder", True))
        self.output_folder_name = ctk.StringVar(value=settings.get("output_folder_name", "output"))
        self.copy_source = ctk.BooleanVar(value=settings.get("copy_source", True))
        self.client_name = ctk.StringVar(value=settings.get("client_name", "client_demo"))
        self.ai_preclean = ctk.BooleanVar(value=settings.get("ai_preclean", True))

        self.api_url = ctk.StringVar(value=settings.get("api_url", "https://openrouter.ai/api/v1/chat/completions"))
        self.api_key = ctk.StringVar(value=settings.get("api_key", ""))
        self.model_name = ctk.StringVar(value=settings.get("model_name", "openrouter/free"))
        self.api_status_text = ctk.StringVar(value="● API не проверено")

        self._build_ui()
        self.after(120, self._drain_queue)

    def _build_ui(self) -> None:
        self.grid_columnconfigure(0, weight=0)
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)

        self.sidebar = ctk.CTkScrollableFrame(self, width=370, fg_color=CARD_ALT, corner_radius=0, border_width=1, border_color=BORDER)
        self.sidebar.grid(row=0, column=0, sticky="nsew")

        self.main = ctk.CTkFrame(self, fg_color=APP_BG, corner_radius=0)
        self.main.grid(row=0, column=1, sticky="nsew")
        self.main.grid_columnconfigure(0, weight=1)
        self.main.grid_rowconfigure(2, weight=1)
        self.main.grid_rowconfigure(3, weight=1)

        self._build_sidebar()
        self._build_header()
        self._build_file_panel()
        self._build_bottom_panels()

    def _build_sidebar(self) -> None:
        top = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        top.pack(fill="x", padx=20, pady=(20, 10))
        ctk.CTkLabel(top, text="Doc-to-Content", font=ctk.CTkFont(size=28, weight="bold"), text_color=TEXT).pack(anchor="w")
        ctk.CTkLabel(top, text="GUI + CLI. Ядро уже можно оборачивать в бота, вебку или автоматику.", text_color=TEXT_DIM, wraplength=300, justify="left").pack(anchor="w", pady=(8, 0))

        ctk.CTkButton(self.sidebar, text="Конвертировать сейчас", command=self.start_conversion, height=48, corner_radius=14, fg_color=ACCENT, hover_color=ACCENT_HOVER, font=ctk.CTkFont(size=15, weight="bold")).pack(fill="x", padx=16, pady=(8, 16))

        self._section_label(self.sidebar, "ПАКЕТ РЕЗУЛЬТАТА")
        pack_card = self._card(self.sidebar)
        pack_card.pack(fill="x", padx=16, pady=(0, 12))
        for label, value in [("Single Output", "single"), ("Creator Pack", "creator"), ("Knowledge Pack", "knowledge")]:
            ctk.CTkRadioButton(pack_card, text=label, value=value, variable=self.pack_preset, text_color=TEXT, command=self._toggle_single_output_state).pack(anchor="w", padx=14, pady=6)

        self._section_label(self.sidebar, "SINGLE OUTPUT")
        result_card = self._card(self.sidebar)
        result_card.pack(fill="x", padx=16, pady=(0, 12))
        self.single_output_buttons = []
        for label, value in [("Чистый текст", "text"), ("Markdown", "markdown"), ("Summary", "summary"), ("Статья", "article"), ("Посты", "posts"), ("FAQ", "faq")]:
            rb = ctk.CTkRadioButton(result_card, text=label, value=value, variable=self.single_result_type, text_color=TEXT)
            rb.pack(anchor="w", padx=14, pady=6)
            self.single_output_buttons.append(rb)

        self._section_label(self.sidebar, "ОБРАБОТКА")
        mode_card = self._card(self.sidebar)
        mode_card.pack(fill="x", padx=16, pady=(0, 12))
        for text, value in [("RAW — как есть", "RAW"), ("CLEAN — локальная чистка", "CLEAN"), ("AI — через модель", "AI")]:
            ctk.CTkRadioButton(mode_card, text=text, value=value, variable=self.process_mode, text_color=TEXT).pack(anchor="w", padx=14, pady=8)
        ctk.CTkCheckBox(mode_card, text="Перед AI сначала прогонять CLEAN", variable=self.ai_preclean, text_color=TEXT).pack(anchor="w", padx=14, pady=(6, 12))

        self._section_label(self.sidebar, "УПАКОВКА")
        package_card = self._card(self.sidebar)
        package_card.pack(fill="x", padx=16, pady=(0, 12))
        self._entry(package_card, "Имя клиента / проекта", self.client_name)
        ctk.CTkCheckBox(package_card, text="Сохранять в отдельную папку", variable=self.use_output_folder, text_color=TEXT).pack(anchor="w", padx=14, pady=(12, 10))
        ctk.CTkEntry(package_card, textvariable=self.output_folder_name, fg_color=APP_BG, border_color=BORDER, text_color=TEXT).pack(fill="x", padx=14, pady=(0, 12))
        ctk.CTkCheckBox(package_card, text="Копировать исходник в пакет", variable=self.copy_source, text_color=TEXT).pack(anchor="w", padx=14, pady=(0, 14))

        self._section_label(self.sidebar, "AI НАСТРОЙКИ")
        ai_card = self._card(self.sidebar)
        ai_card.pack(fill="x", padx=16, pady=(0, 12))
        self._entry(ai_card, "API URL", self.api_url)
        self._entry(ai_card, "API Key", self.api_key, show="*")
        self._entry(ai_card, "Model", self.model_name)
        ctk.CTkLabel(ai_card, textvariable=self.api_status_text, text_color=TEXT_DIM, font=ctk.CTkFont(size=12, weight="bold")).pack(anchor="w", padx=14, pady=(4, 10))

        api_btns = ctk.CTkFrame(ai_card, fg_color="transparent")
        api_btns.pack(fill="x", padx=14, pady=(0, 8))
        ctk.CTkButton(api_btns, text="Проверить API", command=self.check_api, height=38, corner_radius=12, fg_color="#1b2943", hover_color="#243657", border_width=1, border_color=BORDER).pack(side="left", fill="x", expand=True, padx=(0, 6))
        ctk.CTkButton(api_btns, text="Вставить OpenRouter", command=self.fill_openrouter_defaults, height=38, corner_radius=12, fg_color="#1b2943", hover_color="#243657", border_width=1, border_color=BORDER).pack(side="left", fill="x", expand=True, padx=(6, 0))

        bottom_btns = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        bottom_btns.pack(fill="x", padx=16, pady=(0, 18))
        ctk.CTkButton(bottom_btns, text="Открыть output", command=self.open_output_folder, height=42, corner_radius=14, fg_color="#1b2943", hover_color="#243657", border_width=1, border_color=BORDER).pack(fill="x", pady=(0, 10))
        ctk.CTkButton(bottom_btns, text="Сохранить настройки", command=self.save_settings, height=42, corner_radius=14, fg_color="#1b2943", hover_color="#243657", border_width=1, border_color=BORDER).pack(fill="x", pady=(0, 10))
        ctk.CTkButton(bottom_btns, text="Очистить список", command=self.clear_files, height=42, corner_radius=14, fg_color="#1b2943", hover_color="#243657", border_width=1, border_color=BORDER).pack(fill="x")

        self._toggle_single_output_state()

    def _build_header(self) -> None:
        header = ctk.CTkFrame(self.main, fg_color="transparent")
        header.grid(row=0, column=0, sticky="ew", padx=26, pady=(22, 10))
        header.grid_columnconfigure(0, weight=1)
        left = ctk.CTkFrame(header, fg_color="transparent")
        left.grid(row=0, column=0, sticky="w")
        ctk.CTkLabel(left, text="Файлы", font=ctk.CTkFont(size=30, weight="bold"), text_color=TEXT).pack(anchor="w")
        ctk.CTkLabel(left, text="Теперь ядро можно запускать и без окна — через CLI.", text_color=TEXT_DIM).pack(anchor="w", pady=(6, 0))

        actions = ctk.CTkFrame(header, fg_color="transparent")
        actions.grid(row=0, column=1, sticky="e")
        ctk.CTkButton(actions, text="Добавить файлы", command=self.add_files, width=150, height=40, corner_radius=14, fg_color="#1b2943", hover_color="#243657", border_width=1, border_color=BORDER).pack(side="left", padx=(0, 10))
        ctk.CTkButton(actions, text="Добавить папку", command=self.add_folder, width=150, height=40, corner_radius=14, fg_color="#1b2943", hover_color="#243657", border_width=1, border_color=BORDER).pack(side="left", padx=(0, 10))
        ctk.CTkButton(actions, text="Конвертировать", command=self.start_conversion, width=170, height=40, corner_radius=14, fg_color=ACCENT, hover_color=ACCENT_HOVER).pack(side="left")

        progress_wrap = ctk.CTkFrame(self.main, fg_color=CARD_BG, corner_radius=18, border_width=1, border_color=BORDER, height=76)
        progress_wrap.grid(row=1, column=0, sticky="ew", padx=26, pady=(0, 14))
        progress_wrap.grid_columnconfigure(0, weight=1)
        self.status_label = ctk.CTkLabel(progress_wrap, text="Готов к работе.", text_color=TEXT, anchor="w", font=ctk.CTkFont(size=14, weight="bold"))
        self.status_label.grid(row=0, column=0, sticky="ew", padx=18, pady=(14, 8))
        self.progress = ctk.CTkProgressBar(progress_wrap, progress_color=ACCENT, fg_color="#22314f", corner_radius=999)
        self.progress.grid(row=1, column=0, sticky="ew", padx=18, pady=(0, 14))
        self.progress.set(0)

    def _build_file_panel(self) -> None:
        zone = ctk.CTkFrame(self.main, fg_color=CARD_BG, corner_radius=22, border_width=1, border_color=BORDER)
        zone.grid(row=2, column=0, sticky="nsew", padx=26, pady=(0, 14))
        zone.grid_columnconfigure(0, weight=1)
        zone.grid_rowconfigure(1, weight=1)

        top = ctk.CTkFrame(zone, fg_color="transparent")
        top.grid(row=0, column=0, sticky="ew", padx=18, pady=(18, 10))
        top.grid_columnconfigure(0, weight=1)
        ctk.CTkLabel(top, text="Список файлов", text_color=TEXT, font=ctk.CTkFont(size=18, weight="bold")).grid(row=0, column=0, sticky="w")
        self.file_count_label = ctk.CTkLabel(top, text="0 файлов", text_color=TEXT_DIM)
        self.file_count_label.grid(row=0, column=1, sticky="e")

        self.file_box = ctk.CTkTextbox(zone, fg_color=APP_BG, border_width=0, text_color=TEXT, wrap="none", corner_radius=16, font=("Consolas", 13))
        self.file_box.grid(row=1, column=0, sticky="nsew", padx=18, pady=(0, 18))
        self.file_box.insert("1.0", "Пока пусто. Добавь файлы или папку через кнопки сверху.")
        self.file_box.configure(state="disabled")

    def _build_bottom_panels(self) -> None:
        panel_wrap = ctk.CTkFrame(self.main, fg_color="transparent")
        panel_wrap.grid(row=3, column=0, sticky="nsew", padx=26, pady=(0, 22))
        panel_wrap.grid_columnconfigure(0, weight=1)
        panel_wrap.grid_columnconfigure(1, weight=1)
        panel_wrap.grid_rowconfigure(0, weight=1)

        preview_card = ctk.CTkFrame(panel_wrap, fg_color=CARD_BG, corner_radius=22, border_width=1, border_color=BORDER)
        preview_card.grid(row=0, column=0, sticky="nsew", padx=(0, 7))
        ctk.CTkLabel(preview_card, text="Превью первого файла", text_color=TEXT, font=ctk.CTkFont(size=18, weight="bold")).pack(anchor="w", padx=18, pady=(16, 10))
        self.preview_box = ctk.CTkTextbox(preview_card, fg_color=APP_BG, border_width=0, text_color=TEXT_DIM, font=("Consolas", 12))
        self.preview_box.pack(fill="both", expand=True, padx=18, pady=(0, 18))

        log_card = ctk.CTkFrame(panel_wrap, fg_color=CARD_BG, corner_radius=22, border_width=1, border_color=BORDER)
        log_card.grid(row=0, column=1, sticky="nsew", padx=(7, 0))
        ctk.CTkLabel(log_card, text="Лог", text_color=TEXT, font=ctk.CTkFont(size=18, weight="bold")).pack(anchor="w", padx=18, pady=(16, 10))
        self.log_box = ctk.CTkTextbox(log_card, fg_color=APP_BG, border_width=0, text_color=TEXT_DIM, font=("Consolas", 12))
        self.log_box.pack(fill="both", expand=True, padx=18, pady=(0, 18))

    def _section_label(self, parent, text: str) -> None:
        ctk.CTkLabel(parent, text=text, text_color=TEXT_DIM, font=ctk.CTkFont(size=11, weight="bold")).pack(anchor="w", padx=18, pady=(2, 8))

    def _card(self, parent):
        return ctk.CTkFrame(parent, fg_color=CARD_BG, corner_radius=18, border_width=1, border_color=BORDER)

    def _entry(self, parent, label: str, var, show: str | None = None) -> None:
        ctk.CTkLabel(parent, text=label, text_color=TEXT_DIM).pack(anchor="w", padx=14, pady=(12, 6))
        ctk.CTkEntry(parent, textvariable=var, show=show, fg_color=APP_BG, border_color=BORDER, text_color=TEXT, height=38).pack(fill="x", padx=14, pady=(0, 2))

    def _toggle_single_output_state(self) -> None:
        enabled = self.pack_preset.get() == "single"
        state = "normal" if enabled else "disabled"
        for rb in getattr(self, "single_output_buttons", []):
            rb.configure(state=state)

    def log(self, text: str) -> None:
        self.log_box.insert("end", text + "\n")
        self.log_box.see("end")

    def set_status(self, text: str) -> None:
        self.status_label.configure(text=text)

    def set_progress(self, value: float) -> None:
        self.progress.set(max(0.0, min(1.0, value)))

    def fill_openrouter_defaults(self) -> None:
        self.api_url.set("https://openrouter.ai/api/v1/chat/completions")
        self.model_name.set("openrouter/free")
        self.api_status_text.set("● Подставлены дефолты OpenRouter")
        self.log("Подставлены значения OpenRouter по умолчанию.")

    def check_api(self) -> None:
        thread = threading.Thread(target=self._check_api_worker, daemon=True)
        thread.start()

    def _check_api_worker(self) -> None:
        self.queue.put(("api_status", "● Проверяю API..."))
        models_to_try = [self.model_name.get().strip()] if self.model_name.get().strip() else []
        for fallback in FREE_MODELS:
            if fallback not in models_to_try:
                models_to_try.append(fallback)
        for candidate in models_to_try:
            try:
                response = requests.post(
                    self.api_url.get().strip(),
                    headers=_build_headers(self.api_key.get()),
                    json={
                        "model": candidate,
                        "messages": [{"role": "user", "content": "ping"}],
                        "max_tokens": 5,
                        "temperature": 0,
                    },
                    timeout=20,
                )
                response.raise_for_status()
                self.queue.put(("api_status", f"● API подключен ({candidate})"))
                self.queue.put(("log", f"API OK: {candidate}"))
                if candidate != self.model_name.get().strip():
                    self.model_name.set(candidate)
                return
            except Exception as exc:
                self.queue.put(("log", f"API ERROR on {candidate}: {exc}"))
        self.queue.put(("api_status", "● Ошибка API"))

    def save_settings(self) -> None:
        data = {
            "process_mode": self.process_mode.get(),
            "single_result_type": self.single_result_type.get(),
            "pack_preset": self.pack_preset.get(),
            "use_output_folder": self.use_output_folder.get(),
            "output_folder_name": self.output_folder_name.get(),
            "copy_source": self.copy_source.get(),
            "client_name": self.client_name.get(),
            "ai_preclean": self.ai_preclean.get(),
            "api_url": self.api_url.get(),
            "api_key": self.api_key.get(),
            "model_name": self.model_name.get(),
        }
        write_json(SETTINGS_PATH, data)
        self.log(f"Настройки сохранены: {SETTINGS_PATH}")

    def add_files(self) -> None:
        paths = filedialog.askopenfilenames(title="Выбери файлы", filetypes=[("Supported", "*.pdf *.docx *.pptx *.html *.htm"), ("All files", "*.*")])
        self._append_paths([Path(p) for p in paths])

    def add_folder(self) -> None:
        folder = filedialog.askdirectory(title="Выбери папку")
        if not folder:
            return
        paths = collect_input_paths(Path(folder))
        self._append_paths(paths)

    def _append_paths(self, paths: list[Path]) -> None:
        added = 0
        for p in paths:
            if not p.exists() or p.suffix.lower() not in SUPPORTED_INPUTS:
                continue
            if p not in self.files:
                self.files.append(p)
                added += 1
        self.files.sort(key=lambda x: x.name.lower())
        self._refresh_file_list()
        if added:
            self._preview_first_file()
            self.log(f"Добавлено файлов: {added}")

    def clear_files(self) -> None:
        self.files.clear()
        self._refresh_file_list()
        self.preview_box.delete("1.0", "end")
        self.log("Список файлов очищен.")
        self.set_status("Готов к работе.")
        self.set_progress(0)

    def _refresh_file_list(self) -> None:
        self.file_box.configure(state="normal")
        self.file_box.delete("1.0", "end")
        if not self.files:
            self.file_box.insert("1.0", "Пока пусто. Добавь файлы или папку через кнопки сверху.")
        else:
            for i, p in enumerate(self.files, start=1):
                self.file_box.insert("end", f"[{i:02}] {p.name}\n    {p}\n\n")
        self.file_box.configure(state="disabled")
        self.file_count_label.configure(text=f"{len(self.files)} файлов")

    def _preview_first_file(self) -> None:
        if not self.files:
            return
        try:
            raw = extract_raw_text(self.files[0], ".md")
            preview = raw[:3000]
            self.preview_box.delete("1.0", "end")
            self.preview_box.insert("1.0", preview)
        except Exception as exc:
            self.preview_box.delete("1.0", "end")
            self.preview_box.insert("1.0", f"Не удалось показать превью:\n{exc}")

    def open_output_folder(self) -> None:
        if not self.files:
            messagebox.showwarning("Нет файлов", "Сначала добавь хотя бы один файл.")
            return
        base = self.files[0].parent / self.output_folder_name.get().strip()
        if not base.exists():
            base.mkdir(parents=True, exist_ok=True)
        try:
            import os
            os.startfile(str(base))  # type: ignore[attr-defined]
        except Exception as exc:
            messagebox.showerror("Ошибка", f"Не удалось открыть папку: {exc}")

    def start_conversion(self) -> None:
        if not self.files:
            self.log("Нет файлов для обработки.")
            self.set_status("Добавь хотя бы один файл.")
            return
        self.save_settings()
        thread = threading.Thread(target=self._worker, daemon=True)
        thread.start()

    def _worker(self) -> None:
        total = len(self.files)
        preset = self.pack_preset.get()
        jobs = [(self.single_result_type.get(), ".md" if self.single_result_type.get() != "text" else ".txt")] if preset == "single" else PACK_PRESETS[preset]

        service = CoreService(
            process_mode=self.process_mode.get(),
            client_name=self.client_name.get(),
            output_folder_name=self.output_folder_name.get(),
            use_output_folder=self.use_output_folder.get(),
            copy_source=self.copy_source.get(),
            api_url=self.api_url.get(),
            api_key=self.api_key.get(),
            model_name=self.model_name.get(),
            ai_preclean=self.ai_preclean.get(),
            logger=lambda msg: self.queue.put(("log", msg)),
        )

        self.queue.put(("status", f"Погнали. Файлов: {total}"))
        self.queue.put(("progress", 0.0))
        for index, src in enumerate(self.files, start=1):
            try:
                client_dir = service.process_file(src, jobs)
                self.queue.put(("result", JobResult(src, client_dir, True, "ok")))
            except Exception as exc:
                self.queue.put(("result", JobResult(src, None, False, str(exc))))
                self.queue.put(("log", f"Ошибка: {src.name} -> {exc}"))
            self.queue.put(("progress", index / total))
            self.queue.put(("status", f"Обработано {index}/{total}"))
        self.queue.put(("status", "Готово. Пакеты собраны."))
        self.queue.put(("log", f"История сохранена в {HISTORY_PATH}"))
        self.queue.put(("log", "Готово всё."))

    def _drain_queue(self) -> None:
        try:
            while True:
                kind, payload = self.queue.get_nowait()
                if kind == "log":
                    self.log(str(payload))
                elif kind == "status":
                    self.set_status(str(payload))
                elif kind == "progress":
                    self.set_progress(float(payload))
                elif kind == "api_status":
                    self.api_status_text.set(str(payload))
                elif kind == "result":
                    result: JobResult = payload
                    if result.ok and result.target:
                        self.log(f"OK  {result.source.name} -> {result.target}")
                    else:
                        self.log(f"ERR {result.source.name} -> {result.message}")
        except queue.Empty:
            pass
        finally:
            self.after(120, self._drain_queue)


# -------------------- CLI --------------------

def build_jobs_from_args(preset: str, result: str) -> list[tuple[str, str]]:
    if preset == "single":
        return [(result, ".md" if result != "text" else ".txt")]
    return PACK_PRESETS[preset]


def run_cli(args: argparse.Namespace) -> int:
    input_path = Path(args.input)
    files = collect_input_paths(input_path)
    if not files:
        print(f"Не найдено поддерживаемых файлов: {input_path}")
        return 1

    jobs = build_jobs_from_args(args.preset, args.result)
    service = CoreService(
        process_mode=args.mode,
        client_name=args.client,
        output_folder_name=args.output_dir,
        use_output_folder=not args.no_output_folder,
        copy_source=not args.no_copy_source,
        api_url=args.api_url,
        api_key=args.api_key,
        model_name=args.model,
        ai_preclean=not args.no_preclean,
        logger=lambda msg: print(msg),
    )

    created_dirs: list[Path] = []
    for src in files:
        created_dirs.append(service.process_file(src, jobs))

    if args.zip:
        for directory in created_dirs:
            zip_path = directory.with_suffix(".zip")
            zip_directory(directory, zip_path)
            print(f"ZIP: {zip_path}")

    print("Готово.")
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Doc-to-Content v10")
    parser.add_argument("--cli", action="store_true", help="Запустить без GUI")
    parser.add_argument("--input", help="Файл или папка для обработки")
    parser.add_argument("--client", default="client_demo", help="Имя клиента или проекта")
    parser.add_argument("--preset", default="single", choices=["single", "creator", "knowledge"])
    parser.add_argument("--result", default="summary", choices=["text", "markdown", "summary", "article", "posts", "faq"])
    parser.add_argument("--mode", default="AI", choices=["RAW", "CLEAN", "AI"])
    parser.add_argument("--api-url", default="https://openrouter.ai/api/v1/chat/completions")
    parser.add_argument("--api-key", default="")
    parser.add_argument("--model", default="openrouter/free")
    parser.add_argument("--output-dir", default="output")
    parser.add_argument("--no-output-folder", action="store_true")
    parser.add_argument("--no-copy-source", action="store_true")
    parser.add_argument("--no-preclean", action="store_true")
    parser.add_argument("--zip", action="store_true")
    return parser


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()

    if args.cli:
        if not args.input:
            parser.error("Для --cli нужен --input")
        return run_cli(args)

    app = App()
    app.mainloop()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
